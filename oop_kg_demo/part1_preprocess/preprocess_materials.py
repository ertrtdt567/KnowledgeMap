"""
Part 1: preprocessing course materials for the OOP knowledge graph demo.

Input materials such as PPTX/PDF/TXT/MD/JSON exercises are converted into a
unified clean_chunks.json file. This step does not build graph entities or
relations; it prepares clean, traceable chunks for later extraction.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any


OOP_KEYWORDS = {
    "面向对象", "类", "对象", "属性", "方法", "成员变量", "构造方法", "构造函数",
    "封装", "继承", "多态", "抽象", "抽象类", "接口", "重载", "重写",
    "父类", "子类", "this", "super", "class", "object", "method", "field",
    "constructor", "encapsulation", "inheritance", "polymorphism", "abstract",
    "interface", "extends", "implements", "override", "overload", "private",
    "protected", "public",
}

STRONG_OOP_KEYWORDS = {
    "面向对象", "构造方法", "构造函数", "封装", "继承", "多态", "抽象类",
    "重载", "重写", "父类", "子类", "this", "super", "extends", "implements",
    "override", "overload", "inheritance", "polymorphism", "encapsulation",
}

WEAK_OOP_KEYWORDS = {
    "类", "对象", "属性", "方法", "成员变量", "接口", "抽象", "class", "object",
    "method", "field", "constructor", "abstract", "interface", "private",
    "protected", "public",
}

BACKGROUND_TOPIC_KEYWORDS = {
    "课程内容", "课程安排", "课程目标", "参考教材", "参考文献", "java概述",
    "java的历史", "历史", "发展历程", "版本", "版本号", "jdk", "jre", "ide",
    "eclipse", "安装", "环境配置", "sun公司", "oracle", "james gosling", "平台",
    "虚拟机", "字节码", "垃圾收集器", "hello world", "helloworld",
    "public static void main", "程序入口点", "源文件扩展名", "输入输出", "线程",
    "网络程序设计", "gui",
}

NOISE_LINE_PATTERNS = [
    r"^\s*\d+\s*$", r"^\s*第\s*\d+\s*页\s*$", r"^\s*page\s*\d+\s*$",
    r"^\s*目录\s*$", r"^\s*contents?\s*$", r"^\s*thank\s+you\s*$",
    r"^\s*谢谢观看\s*$", r"^\s*参考文献\s*$", r"^\s*references?\s*$",
]

ROLE_KEYWORDS = {
    "code_example": {"class ", "public class", "extends ", "implements ", "new ", "void ", "{", "}", ";"},
    "exercise": {"选择题", "填空题", "判断题", "编程题", "下列", "正确的是", "错误的是", "答案", "解析"},
    "syntax_rule": {"语法", "关键字", "使用", "格式", "extends", "implements", "class", "interface"},
    "concept_explanation": {"定义", "概念", "特点", "作用", "表示", "体现", "机制"},
}

NORMALIZATION_MAP = {
    "OOP": "面向对象编程", "Object-Oriented Programming": "面向对象编程",
    "object-oriented programming": "面向对象编程", "class": "类", "object": "对象",
    "method": "方法", "field": "属性", "member variable": "属性", "constructor": "构造方法",
    "inheritance": "继承", "polymorphism": "多态", "encapsulation": "封装", "interface": "接口",
}


@dataclass
class RawDocument:
    source_file: str
    source_type: str
    raw_text: str
    page: int | None = None
    title: str = ""
    question_id: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class CleanChunk:
    chunk_id: str
    source_file: str
    source_type: str
    language: str
    page: int | None
    chapter: str
    section: str
    content: str
    normalized_content: str
    keywords: list[str]
    material_role: str
    evidence_location: str
    confidence_hint: float
    metadata: dict[str, Any] = field(default_factory=dict)


class MaterialLoader:
    """Read supported files and convert them into RawDocument records."""

    def load_folder(self, input_dir: Path) -> list[RawDocument]:
        docs: list[RawDocument] = []
        for file_path in sorted(input_dir.rglob("*")):
            if not file_path.is_file():
                continue
            suffix = file_path.suffix.lower()
            if suffix == ".pptx":
                docs.extend(self.load_pptx(file_path))
            elif suffix == ".pdf":
                docs.extend(self.load_pdf(file_path))
            elif suffix in {".txt", ".md"}:
                docs.extend(self.load_text(file_path))
            elif suffix == ".json":
                docs.extend(self.load_json(file_path))
        return docs

    def load_pptx(self, file_path: Path) -> list[RawDocument]:
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("Please install python-pptx to read PPTX files.") from exc
        docs: list[RawDocument] = []
        for page, slide in enumerate(Presentation(str(file_path)).slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            raw_text = "\n".join(texts).strip()
            if raw_text:
                docs.append(RawDocument(file_path.name, "ppt", raw_text, page, self._guess_title(raw_text)))
        return docs

    def load_pdf(self, file_path: Path) -> list[RawDocument]:
        try:
            from pypdf import PdfReader
        except ImportError:
            from PyPDF2 import PdfReader  # type: ignore
        docs: list[RawDocument] = []
        for page_index, page in enumerate(PdfReader(str(file_path)).pages, start=1):
            raw_text = (page.extract_text() or "").strip()
            if raw_text:
                docs.append(RawDocument(file_path.name, "pdf", raw_text, page_index, self._guess_title(raw_text)))
        return docs

    def load_text(self, file_path: Path) -> list[RawDocument]:
        raw_text = file_path.read_text(encoding="utf-8", errors="ignore").strip()
        return [RawDocument(file_path.name, "text", raw_text, title=self._guess_title(raw_text))] if raw_text else []

    def load_json(self, file_path: Path) -> list[RawDocument]:
        data = json.loads(file_path.read_text(encoding="utf-8"))
        records = data if isinstance(data, list) else data.get("questions", [])
        docs: list[RawDocument] = []
        for index, item in enumerate(records, start=1):
            if not isinstance(item, dict):
                continue
            qid = str(item.get("id") or item.get("question_id") or index)
            raw_text = self._format_exercise(
                str(item.get("question") or item.get("stem") or ""),
                item.get("options") or [],
                str(item.get("answer") or ""),
                str(item.get("analysis") or item.get("explanation") or ""),
            )
            if raw_text.strip():
                docs.append(RawDocument(file_path.name, "exercise", raw_text, question_id=qid, title=f"习题 {qid}"))
        return docs

    @staticmethod
    def _guess_title(text: str) -> str:
        return next((line.strip()[:60] for line in text.splitlines() if line.strip()), "")

    @staticmethod
    def _format_exercise(question: str, options: Any, answer: str, analysis: str) -> str:
        lines = [f"题干：{question}"] if question else []
        if isinstance(options, list):
            lines.extend(str(x) for x in options)
        elif isinstance(options, dict):
            lines.extend(f"{k}. {v}" for k, v in options.items())
        if answer:
            lines.append(f"答案：{answer}")
        if analysis:
            lines.append(f"解析：{analysis}")
        return "\n".join(lines)


class TextCleaner:
    def clean(self, text: str) -> str:
        text = text.replace("\u3000", " ")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
        lines = [re.sub(r"\s+", " ", line.strip()) for line in text.splitlines()]
        lines = [line for line in lines if line and not self._is_noise_line(line)]
        return re.sub(r"\n{3,}", "\n\n", "\n".join(lines)).strip()

    @staticmethod
    def _is_noise_line(line: str) -> bool:
        return any(re.match(pattern, line, flags=re.IGNORECASE) for pattern in NOISE_LINE_PATTERNS)


class DomainFilter:
    def matched_keywords(self, text: str) -> list[str]:
        lower_text = text.lower()
        return sorted({k for k in OOP_KEYWORDS if k.lower() in lower_text}, key=lambda x: (len(x), x))

    def keep(self, doc: RawDocument, cleaned_text: str) -> bool:
        return self._relevance_score(cleaned_text) >= 3

    def keep_chunk(self, doc: RawDocument, chunk_text: str) -> bool:
        score = self._relevance_score(chunk_text)
        return score >= (2 if doc.source_type == "exercise" else 3)

    def _relevance_score(self, text: str) -> int:
        lower_text = text.lower()
        strong = [k for k in STRONG_OOP_KEYWORDS if k.lower() in lower_text]
        weak = [k for k in WEAK_OOP_KEYWORDS if k.lower() in lower_text]
        score = len(strong) * 3 + min(len(weak), 3)
        if re.search(r"\b(class|interface|extends|implements|abstract\s+class)\b", text, flags=re.IGNORECASE):
            score += 2
        if any(k.lower() in lower_text for k in BACKGROUND_TOPIC_KEYWORDS):
            score -= 5
        return score


class ChunkSplitter:
    def split(self, doc: RawDocument, cleaned_text: str) -> list[str]:
        if doc.source_type == "exercise":
            return [cleaned_text]
        blocks = re.split(r"\n\s*\n", cleaned_text)
        chunks: list[str] = []
        for block in blocks:
            lines = [line.strip() for line in block.splitlines() if line.strip()]
            current: list[str] = []
            for line in lines:
                is_bullet = bool(re.match(r"^([\-*•]|\d+[.)、]|[A-D][.)、])\s*", line))
                if is_bullet and current:
                    chunks.append(" ".join(current))
                    current = [line]
                else:
                    current.append(line)
            if current:
                chunks.extend(self._split_long(" ".join(current)))
        return [c.strip() for c in chunks if len(c.strip()) >= 8]

    @staticmethod
    def _split_long(text: str) -> list[str]:
        if len(text) <= 420:
            return [text]
        parts = re.split(r"(?<=[。！？.!?])\s+", text)
        chunks: list[str] = []
        current = ""
        for part in parts:
            if len(current) + len(part) > 420 and current:
                chunks.append(current)
                current = part
            else:
                current = f"{current} {part}".strip()
        if current:
            chunks.append(current)
        return chunks


class ChunkLabeler:
    def label(self, text: str, doc: RawDocument, keywords: list[str]) -> dict[str, Any]:
        return {
            "language": self._detect_language(text, doc.source_file),
            "material_role": self._detect_role(text, doc.source_type),
            "chapter": self._detect_chapter(doc.title, text),
            "section": doc.title or "、".join(keywords[:3]) or "未标注小节",
            "normalized_content": self._normalize_terms(text),
            "confidence_hint": self._confidence_hint(text, keywords),
        }

    @staticmethod
    def _detect_language(text: str, source_file: str) -> str:
        combined = f"{source_file}\n{text}".lower()
        if "java" in combined or "extends" in combined or "implements" in combined:
            return "Java"
        if "python" in combined or "__init__" in combined or "self." in combined:
            return "Python"
        if "c++" in combined or "virtual" in combined or "::" in combined:
            return "C++"
        return "通用OOP"

    @staticmethod
    def _detect_role(text: str, source_type: str) -> str:
        if source_type == "exercise":
            return "exercise"
        lower = text.lower()
        scores = {role: sum(1 for k in kws if k.lower() in lower) for role, kws in ROLE_KEYWORDS.items()}
        role = max(scores, key=scores.get)
        return role if scores[role] > 0 else "concept_explanation"

    @staticmethod
    def _detect_chapter(title: str, text: str) -> str:
        combined = f"{title} {text}".lower()
        if "继承" in combined or "extends" in combined:
            return "面向对象编程-继承"
        if "多态" in combined or "polymorphism" in combined:
            return "面向对象编程-多态"
        if "接口" in combined or "interface" in combined:
            return "面向对象编程-接口"
        if "封装" in combined or "private" in combined or "protected" in combined:
            return "面向对象编程-封装"
        if "构造" in combined or "constructor" in combined:
            return "面向对象编程-构造方法"
        if "类" in combined or "对象" in combined or "class" in combined:
            return "面向对象编程-类与对象"
        return "面向对象编程"

    @staticmethod
    def _normalize_terms(text: str) -> str:
        normalized = text
        for old, new in NORMALIZATION_MAP.items():
            normalized = re.sub(re.escape(old), new, normalized, flags=re.IGNORECASE)
        return normalized

    @staticmethod
    def _confidence_hint(text: str, keywords: list[str]) -> float:
        score = 0.55 + (0.15 if len(keywords) >= 2 else 0) + (0.1 if len(text) >= 25 else 0)
        if any(mark in text for mark in ["class", "extends", "implements", "接口", "继承", "多态", "封装"]):
            score += 0.15
        return min(score, 0.95)


class MaterialPreprocessor:
    def __init__(self) -> None:
        self.loader = MaterialLoader()
        self.cleaner = TextCleaner()
        self.filter = DomainFilter()
        self.splitter = ChunkSplitter()
        self.labeler = ChunkLabeler()

    def run(self, input_dir: Path) -> list[CleanChunk]:
        chunks: list[CleanChunk] = []
        for doc in self.loader.load_folder(input_dir):
            cleaned = self.cleaner.clean(doc.raw_text)
            if not cleaned or not self.filter.keep(doc, cleaned):
                continue
            for index, content in enumerate(self.splitter.split(doc, cleaned), start=1):
                keywords = self.filter.matched_keywords(content)
                if not keywords or not self.filter.keep_chunk(doc, content):
                    continue
                labels = self.labeler.label(content, doc, keywords)
                chunks.append(CleanChunk(
                    chunk_id=self._chunk_id(doc, index, content),
                    source_file=doc.source_file,
                    source_type=doc.source_type,
                    language=labels["language"],
                    page=doc.page,
                    chapter=labels["chapter"],
                    section=labels["section"],
                    content=content,
                    normalized_content=labels["normalized_content"],
                    keywords=keywords,
                    material_role=labels["material_role"],
                    evidence_location=self._evidence_location(doc),
                    confidence_hint=labels["confidence_hint"],
                    metadata=doc.metadata,
                ))
        return chunks

    @staticmethod
    def _chunk_id(doc: RawDocument, index: int, content: str) -> str:
        location = doc.question_id or doc.page or "doc"
        digest = hashlib.md5(f"{doc.source_file}-{location}-{index}-{content}".encode("utf-8")).hexdigest()[:8]
        stem = re.sub(r"[^a-zA-Z0-9]+", "_", Path(doc.source_file).stem).strip("_").lower()
        return f"{stem}_{location}_{index}_{digest}"

    @staticmethod
    def _evidence_location(doc: RawDocument) -> str:
        if doc.page is not None:
            return f"第{doc.page}页"
        if doc.question_id:
            return f"第{doc.question_id}题"
        return "全文"


def write_chunks(chunks: list[CleanChunk], output_path: Path) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(json.dumps([c.__dict__ for c in chunks], ensure_ascii=False, indent=2), encoding="utf-8")


def parse_args(argv: list[str]) -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Preprocess OOP course materials into clean_chunks.json")
    parser.add_argument("--input", required=True, help="Folder containing PPTX/PDF/TXT/MD/JSON materials")
    parser.add_argument("--output", required=True, help="Output clean_chunks.json path")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    args = parse_args(argv)
    input_dir = Path(args.input)
    if not input_dir.exists():
        print(f"输入目录不存在：{input_dir}", file=sys.stderr)
        return 1
    chunks = MaterialPreprocessor().run(input_dir)
    write_chunks(chunks, Path(args.output))
    print(f"已生成 {len(chunks)} 个知识片段：{args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))

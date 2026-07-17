"""
面向对象知识图谱 Demo 的“材料预处理”程序。

这份代码解决的是知识图谱构建的第一步：

    PPT / PDF / TXT / JSON 习题
    → 读取文件内容
    → 清洗无关文字
    → 按知识点切成小片段
    → 判断是否和面向对象编程相关
    → 输出统一格式的 clean_chunks.json

注意：
1. 这里的输出还不是最终知识图谱。
2. clean_chunks.json 是中间数据，后面会继续交给实体抽取、关系抽取等模块。
3. 这份代码目前主要靠规则处理，后续可以把 ChunkLabeler 换成 LLM 标注。
"""

from __future__ import annotations

import argparse
import hashlib
import json
import os
import re
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from curriculum_catalog import CurriculumCatalog, DEFAULT_CATALOG


# OOP_KEYWORDS 用来记录所有可能和“面向对象编程”相关的词。
# 注意：它只负责“记录命中的关键词”，真正是否保留由 DomainFilter 里的评分规则决定。
OOP_KEYWORDS = {
    "面向对象",
    "类",
    "对象",
    "属性",
    "方法",
    "成员变量",
    "构造方法",
    "构造函数",
    "封装",
    "继承",
    "多态",
    "抽象",
    "抽象类",
    "接口",
    "重载",
    "重写",
    "覆盖",
    "父类",
    "子类",
    "基类",
    "派生类",
    "this",
    "super",
    "class",
    "object",
    "method",
    "field",
    "constructor",
    "encapsulation",
    "inheritance",
    "polymorphism",
    "abstract",
    "interface",
    "extends",
    "implements",
    "override",
    "overload",
    "private",
    "protected",
    "public",
}

# STRONG_OOP_KEYWORDS 是强相关词。
# 这些词通常直接指向面向对象核心知识，例如继承、多态、封装、构造方法。
# 命中这些词，说明文本很可能值得保留。
STRONG_OOP_KEYWORDS = {
    "面向对象",
    "构造方法",
    "构造函数",
    "封装",
    "继承",
    "多态",
    "抽象类",
    "重载",
    "重写",
    "父类",
    "子类",
    "基类",
    "派生类",
    "this",
    "super",
    "extends",
    "implements",
    "override",
    "overload",
    "inheritance",
    "polymorphism",
    "encapsulation",
}

# WEAK_OOP_KEYWORDS 是弱相关词。
# 这些词单独出现时不一定代表 OOP 知识，例如“类”可能出现在“分类”“类别”里；
# “接口”也可能是“用户接口”，不一定是 Java interface。
# 所以弱相关词必须组合出现，或者和代码形态一起出现，才会被保留。
WEAK_OOP_KEYWORDS = {
    "类",
    "对象",
    "属性",
    "方法",
    "成员变量",
    "接口",
    "抽象",
    "覆盖",
    "class",
    "object",
    "method",
    "field",
    "constructor",
    "abstract",
    "interface",
    "private",
    "protected",
    "public",
}

# BACKGROUND_TOPIC_KEYWORDS 是偏背景介绍或课程说明的词。
# 如果一页主要在讲 Java 历史、版本、环境配置、课程目录等，即使出现“类/接口”等泛词，
# 也不应该优先进入面向对象知识图谱。
BACKGROUND_TOPIC_KEYWORDS = {
    "课程内容",
    "课程安排",
    "课程目标",
    "参考教材",
    "参考文献",
    "java概述",
    "java的历史",
    "历史",
    "发展历程",
    "版本",
    "版本号",
    "java的特征(1)",
    "简单性",
    "基本语法",
    "语法机制概述",
    "数据类型",
    "表达式",
    "程序流控制",
    "jdk",
    "jre",
    "ide",
    "eclipse",
    "安装",
    "环境配置",
    "sun公司",
    "oracle",
    "james gosling",
    "平台",
    "虚拟机",
    "字节码",
    "垃圾收集器",
    "hello world",
    "helloworld",
    "public static void main",
    "程序入口点",
    "源文件扩展名",
    "java应用程序举例",
    "java程序的编写",
    "找出并修改下面代码的错误",
    "咨询问题",
    "计算机学院",
    "输入输出",
    "线程",
    "网络程序设计",
    "gui",
}

# NOISE_LINE_PATTERNS 用来识别噪声行。
# 例如页码、目录、Thank you 等内容通常不应该进入知识图谱。
NOISE_LINE_PATTERNS = [
    r"^\s*\d+\s*$",
    r"^\s*第\s*\d+\s*页\s*$",
    r"^\s*page\s*\d+\s*$",
    r"^\s*目录\s*$",
    r"^\s*contents?\s*$",
    r"^\s*thank\s+you\s*$",
    r"^\s*谢谢观看\s*$",
    r"^\s*参考文献\s*$",
    r"^\s*references?\s*$",
]

# ROLE_KEYWORDS 用来粗略判断一个片段的材料类型。
# 例如：包含 class、extends、{} 的内容很可能是代码示例；
# 包含“答案”“解析”的内容很可能是习题。
ROLE_KEYWORDS = {
    "code_example": {
        "class ",
        "public class",
        "extends ",
        "implements ",
        "new ",
        "void ",
        "{",
        "}",
        ";",
        "def ",
        "__init__",
    },
    "exercise": {"选择题", "填空题", "判断题", "编程题", "下列", "正确的是", "错误的是", "答案", "解析"},
    "syntax_rule": {"语法", "关键字", "使用", "格式", "extends", "implements", "class", "interface"},
    "concept_explanation": {"定义", "概念", "特点", "作用", "表示", "体现", "机制"},
}

# NORMALIZATION_MAP 用来做术语归一化。
# 例如：class、Class、类，本质上都可以对齐到“类”这个概念。
NORMALIZATION_MAP = {
    "OOP": "面向对象编程",
    "Object-Oriented Programming": "面向对象编程",
    "object-oriented programming": "面向对象编程",
    "class": "类",
    "object": "对象",
    "method": "方法",
    "field": "属性",
    "member variable": "属性",
    "constructor": "构造方法",
    "inheritance": "继承",
    "polymorphism": "多态",
    "encapsulation": "封装",
    "interface": "接口",
}


@dataclass
class RawDocument:
    """
    原始文档片段。

    文件刚被读取出来时，先统一放到 RawDocument 里。
    这个阶段的文本还没有彻底清洗，也还没有按知识点切分。
    """

    source_file: str
    source_type: str
    raw_text: str
    page: int | None = None
    title: str = ""
    question_id: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)


@dataclass
class CleanChunk:
    """
    清洗后的知识片段。

    CleanChunk 是这个程序最终输出到 clean_chunks.json 的数据结构。
    后面的实体抽取和关系抽取，就是基于这些 chunk 继续处理。
    """

    chunk_id: str
    source_file: str
    source_type: str
    language: str
    page: int | None
    chapter: str
    section: str
    content: str
    normalized_content: str
    keywords: list[str]
    domain_candidates: list[str]
    material_role: str
    evidence_location: str
    confidence_hint: float
    metadata: dict[str, Any] = field(default_factory=dict)


class MaterialLoader:
    """
    文件读取器。

    它负责把不同格式的文件读取成统一的 RawDocument：
    - PPTX：按页读取文字
    - PPT：先自动转成 PPTX，再按页读取文字
    - PDF：按页读取文字
    - TXT/MD：按全文读取
    - JSON：按习题读取
    - DOCX：按段落块读取
    """

    def load_folder(self, input_dir: Path) -> list[RawDocument]:
        """遍历输入文件夹，自动识别文件类型，并调用对应的读取方法。"""
        docs: list[RawDocument] = []
        for file_path in sorted(input_dir.rglob("*")):
            if not file_path.is_file():
                continue
            # 老版 PPT 在读取时会即时转换；这里跳过转换缓存，避免同一课程重复两次。
            if "_converted_pptx" in file_path.parts:
                continue
            suffix = file_path.suffix.lower()
            if suffix == ".pptx":
                docs.extend(self.load_pptx(file_path))
            elif suffix == ".ppt":
                docs.extend(self.load_legacy_ppt(file_path, input_dir))
            elif suffix == ".pdf":
                docs.extend(self.load_pdf(file_path))
            elif suffix in {".txt", ".md"}:
                docs.extend(self.load_text(file_path))
            elif suffix == ".docx":
                docs.extend(self.load_docx(file_path))
            elif suffix == ".json":
                docs.extend(self.load_json(file_path))
        return docs

    def load_legacy_ppt(self, file_path: Path, input_dir: Path) -> list[RawDocument]:
        """
        读取老版 PPT 文件。

        python-pptx 只能直接读取 .pptx，不能直接读取 .ppt。
        所以这里先尝试调用本机 PowerPoint，把 .ppt 自动另存为 .pptx，
        然后再调用 load_pptx 读取。
        """
        converted_dir = input_dir / "_converted_pptx"
        converted_dir.mkdir(parents=True, exist_ok=True)
        converted_path = converted_dir / f"{file_path.stem}.pptx"

        if not converted_path.exists() or converted_path.stat().st_mtime < file_path.stat().st_mtime:
            self._convert_ppt_to_pptx(file_path, converted_path)

        if not converted_path.exists():
            print(f"跳过 {file_path.name}：无法转换为 PPTX。", file=sys.stderr)
            return []

        docs = self.load_pptx(converted_path)
        for doc in docs:
            doc.source_file = file_path.name
            doc.metadata["converted_from"] = str(file_path)
            doc.metadata["converted_pptx"] = str(converted_path)
        return docs

    @staticmethod
    def _convert_ppt_to_pptx(source_path: Path, target_path: Path) -> None:
        """
        调用 PowerPoint COM 接口进行格式转换。

        这个方法依赖 Windows 上安装的 Microsoft PowerPoint。
        如果机器没有 PowerPoint，或者文件损坏，转换会失败。
        """
        script = r"""
$ErrorActionPreference = 'Stop'
$source = [System.IO.Path]::GetFullPath($env:PPT_SOURCE)
$target = [System.IO.Path]::GetFullPath($env:PPT_TARGET)
$powerPoint = $null
$presentation = $null
try {
    $powerPoint = New-Object -ComObject PowerPoint.Application
    $presentation = $powerPoint.Presentations.Open($source, $true, $false, $false)
    $presentation.SaveAs($target, 24)
}
finally {
    if ($presentation -ne $null) {
        $presentation.Close()
    }
    if ($powerPoint -ne $null) {
        $powerPoint.Quit()
    }
}
"""
        env = os.environ.copy()
        env["PPT_SOURCE"] = str(source_path)
        env["PPT_TARGET"] = str(target_path)
        result = subprocess.run(
            [
                "powershell",
                "-NoProfile",
                "-ExecutionPolicy",
                "Bypass",
                "-Command",
                script,
            ],
            capture_output=True,
            text=True,
            env=env,
        )
        if result.returncode != 0:
            print(f"PPT 转换失败：{source_path.name}", file=sys.stderr)
            if result.stderr.strip():
                print(result.stderr.strip(), file=sys.stderr)

    def load_pptx(self, file_path: Path) -> list[RawDocument]:
        """读取 PPTX 文件，把每一页幻灯片转换成一个 RawDocument。"""
        try:
            from pptx import Presentation
        except ImportError as exc:
            raise RuntimeError("缺少 python-pptx，无法读取 PPTX 文件。") from exc

        presentation = Presentation(str(file_path))
        docs: list[RawDocument] = []
        for page, slide in enumerate(presentation.slides, start=1):
            texts: list[str] = []
            for shape in slide.shapes:
                # 普通文本框、标题、项目符号通常都可以通过 shape.text 读取。
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
                # 如果 PPT 中有表格，也把每个单元格中的文字取出来。
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            texts.append(" | ".join(cells))
            raw_text = "\n".join(texts).strip()
            if not raw_text:
                continue
            docs.append(
                RawDocument(
                    source_file=file_path.name,
                    source_type="ppt",
                    page=page,
                    title=self._guess_title(raw_text),
                    raw_text=raw_text,
                )
            )
        return docs

    def load_pdf(self, file_path: Path) -> list[RawDocument]:
        """读取 PDF 文件，把每一页 PDF 转换成一个 RawDocument。"""
        try:
            from pypdf import PdfReader
        except ImportError:
            try:
                from PyPDF2 import PdfReader  # type: ignore
            except ImportError as exc:
                raise RuntimeError("缺少 pypdf 或 PyPDF2，无法读取 PDF 文件。") from exc

        reader = PdfReader(str(file_path))
        docs: list[RawDocument] = []
        for page_index, page in enumerate(reader.pages, start=1):
            # 注意：如果 PDF 是扫描版图片，这里可能提取不到文字，需要后续接 OCR。
            raw_text = page.extract_text() or ""
            raw_text = raw_text.strip()
            if not raw_text:
                continue
            docs.append(
                RawDocument(
                    source_file=file_path.name,
                    source_type="pdf",
                    page=page_index,
                    title=self._guess_title(raw_text),
                    raw_text=raw_text,
                )
            )
        return docs

    def load_text(self, file_path: Path) -> list[RawDocument]:
        """读取 TXT 或 Markdown 文件。"""
        raw_text = file_path.read_text(encoding="utf-8", errors="ignore").strip()
        if not raw_text:
            return []
        return [
            RawDocument(
                source_file=file_path.name,
                source_type="text",
                title=self._guess_title(raw_text),
                raw_text=raw_text,
            )
        ]

    def load_docx(self, file_path: Path) -> list[RawDocument]:
        """读取 Word 文档，并按连续非空段落生成可追溯材料块。"""
        try:
            from docx import Document
        except ImportError as exc:
            raise RuntimeError("缺少 python-docx，无法读取 DOCX 文件。") from exc

        document = Document(str(file_path))
        paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        if not paragraphs:
            return []
        blocks: list[list[str]] = []
        current: list[str] = []
        for paragraph in paragraphs:
            current.append(paragraph)
            if len(current) >= 12 or (len(paragraph) < 60 and re.match(r"^(第[一二三四五六七八九十\d]+[篇章节]|\d+[.、．])", paragraph)):
                blocks.append(current)
                current = []
        if current:
            blocks.append(current)
        return [
            RawDocument(
                source_file=file_path.name,
                source_type="docx",
                page=index,
                title=self._guess_title("\n".join(block)),
                raw_text="\n".join(block),
            )
            for index, block in enumerate(blocks, start=1)
            if block
        ]

    def load_json(self, file_path: Path) -> list[RawDocument]:
        """
        读取 JSON 格式的习题文件。

        支持两种常见形式：
        1. 整个 JSON 本身就是一个习题列表。
        2. JSON 里有一个 questions 字段，里面是习题列表。
        """
        data = json.loads(file_path.read_text(encoding="utf-8"))
        records = data if isinstance(data, list) else data.get("questions", [])
        docs: list[RawDocument] = []
        for index, item in enumerate(records, start=1):
            if not isinstance(item, dict):
                continue
            question_id = str(item.get("id") or item.get("question_id") or index)
            question = str(item.get("question") or item.get("stem") or "")
            options = item.get("options") or []
            answer = str(item.get("answer") or "")
            analysis = str(item.get("analysis") or item.get("explanation") or "")
            raw_text = self._format_exercise(question, options, answer, analysis)
            if not raw_text.strip():
                continue
            docs.append(
                RawDocument(
                    source_file=file_path.name,
                    source_type="exercise",
                    question_id=question_id,
                    title=f"习题 {question_id}",
                    raw_text=raw_text,
                    metadata={"answer": answer, "analysis": analysis, "options": options},
                )
            )
        return docs

    @staticmethod
    def _guess_title(text: str) -> str:
        """用第一行非空文本当作标题，简单但够 Demo 使用。"""
        for line in text.splitlines():
            line = line.strip()
            if line:
                return line[:60]
        return ""

    @staticmethod
    def _format_exercise(question: str, options: Any, answer: str, analysis: str) -> str:
        """把题干、选项、答案、解析拼成一段统一文本。"""
        lines = []
        if question:
            lines.append(f"题干：{question}")
        if isinstance(options, list):
            for option in options:
                lines.append(str(option))
        elif isinstance(options, dict):
            for key, value in options.items():
                lines.append(f"{key}. {value}")
        if answer:
            lines.append(f"答案：{answer}")
        if analysis:
            lines.append(f"解析：{analysis}")
        return "\n".join(lines)


class TextCleaner:
    """
    文本清洗器。

    它负责做基础规则清洗，例如：
    - 去掉多余空格
    - 去掉控制字符
    - 去掉页码、目录等噪声行
    - 合并过多空行
    """

    def clean(self, text: str) -> str:
        """清洗一段原始文本，返回更干净的文本。"""
        text = text.replace("\u3000", " ")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
        lines = [self._clean_line(line) for line in text.splitlines()]
        lines = [line for line in lines if line and not self._is_noise_line(line)]
        text = "\n".join(lines)
        text = re.sub(r"[ \t]{2,}", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    @staticmethod
    def _clean_line(line: str) -> str:
        """清洗单行文本：去掉首尾空格，并把连续空格压缩成一个。"""
        line = line.strip()
        line = re.sub(r"\s+", " ", line)
        return line

    @staticmethod
    def _is_noise_line(line: str) -> bool:
        """判断某一行是不是页码、目录、Thank you 等噪声。"""
        for pattern in NOISE_LINE_PATTERNS:
            if re.match(pattern, line, flags=re.IGNORECASE):
                return True
        return False


class DomainFilter:
    """
    领域过滤器。

    它负责判断文本是否和“面向对象编程”相关。
    不相关的内容，例如课程安排、考试说明、参考资料等，会被过滤掉。

    这里使用“相关性评分”，不是简单关键词命中：
    - 强 OOP 词，例如继承、多态、封装，权重高。
    - 弱 OOP 词，例如类、方法、接口，必须组合出现才可信。
    - Java 历史、版本、安装环境等背景页会扣分。
    """

    def matched_keywords(self, text: str) -> list[str]:
        """返回文本中命中的 OOP 关键词列表。"""
        lower_text = text.lower()
        matches = []
        for keyword in OOP_KEYWORDS:
            if keyword.lower() in lower_text:
                matches.append(keyword)
        return sorted(set(matches), key=lambda item: (len(item), item))

    def keep(self, doc: RawDocument, cleaned_text: str) -> bool:
        """
        判断一页或一道题是否值得进入后续切分。

        注意：这是“页级过滤”。它先粗略判断这一页有没有必要继续处理。
        后面还会对切分后的每个 chunk 再做一次“片段级过滤”。
        """
        return self._relevance_score(cleaned_text) >= 3

    def keep_chunk(self, doc: RawDocument, chunk_text: str) -> bool:
        """
        判断切分后的单个 chunk 是否保留。

        片段级过滤会更严格一些，因为此时文本已经更短了，
        可以更准确地过滤掉“Java 历史”“课程目录”这类误入内容。
        """
        score = self._relevance_score(chunk_text)
        if doc.source_type == "exercise":
            return score >= 2
        return score >= 3

    def _relevance_score(self, text: str) -> int:
        """
        计算 OOP 相关性分数。

        评分逻辑：
        1. 强相关词每个加 3 分。
        2. 弱相关词每个加 1 分，但最多加 3 分，避免“类/方法/接口”泛滥。
        3. 如果文本看起来像 Java 类、接口、继承代码，再加 2 分。
        4. 如果文本明显是历史、版本、安装环境、课程目录等背景内容，扣 5 分。

        保留阈值目前设为 3 分。这个阈值适合 Demo 初期使用，
        后续如果过滤太严或太松，可以只改这里。
        """
        lower_text = text.lower()
        strong_matches = self._matched_from_set(lower_text, STRONG_OOP_KEYWORDS)
        weak_matches = self._matched_from_set(lower_text, WEAK_OOP_KEYWORDS)

        score = len(strong_matches) * 3
        score += min(len(weak_matches), 3)

        if self._has_oop_code_pattern(text):
            score += 2

        if self._looks_like_background_page(lower_text):
            # 背景页有时也会偶然出现“类”“接口”“继承”等词，
            # 例如 Java 特性介绍里提到“不支持多重继承”。
            # 这类内容不是当前 Demo 的核心 OOP 知识，所以扣分要重一点。
            score -= 5

        return score

    @staticmethod
    def _matched_from_set(lower_text: str, keywords: set[str]) -> list[str]:
        """从指定关键词集合里找出命中的词。"""
        return [keyword for keyword in keywords if keyword.lower() in lower_text]

    @staticmethod
    def _has_oop_code_pattern(text: str) -> bool:
        """
        判断文本中是否出现典型 OOP 代码形态。

        例如：
        - class Student
        - interface Flyable
        - class Student extends Person
        - implements Runnable
        """
        patterns = [
            r"\bclass\s+[A-Za-z_]\w*",
            r"\binterface\s+[A-Za-z_]\w*",
            r"\bextends\s+[A-Za-z_]\w*",
            r"\bimplements\s+[A-Za-z_]\w*",
            r"\babstract\s+class\b",
            r"\b(this|super)\.",
        ]
        return any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in patterns)

    @staticmethod
    def _looks_like_background_page(lower_text: str) -> bool:
        """
        判断文本是否更像背景介绍页。

        这类内容可能包含“类”“接口”等泛词，但它不是我们当前 Demo 的核心：
        面向对象范式、类与对象、封装、继承、多态、接口、抽象类。
        """
        return any(keyword.lower() in lower_text for keyword in BACKGROUND_TOPIC_KEYWORDS)


class ProgrammingDomainFilter:
    """编程领域过滤器：保留标准知识树覆盖的课程内容，而非只保留 OOP。"""

    CODE_PATTERNS = [
        r"\b(class|struct|template|namespace|public|private|protected)\b",
        r"\b(if|else|for|while|switch|return|break|continue)\b",
        r"\b(try|catch|throw|except|raise|import|def)\b",
        r"[#{};]|\b(int|double|char|bool|void|string)\b",
    ]
    LANGUAGE_MARKERS = {"java", "python", "c++", "c#", "cpp"}

    def __init__(self, catalog: CurriculumCatalog) -> None:
        self.catalog = catalog

    def matched_keywords(self, text: str) -> list[str]:
        return [str(node.get("name", "")) for node in self.catalog.match_text(text) if node.get("name")]

    def domain_candidates(self, text: str) -> list[str]:
        return self.catalog.domain_candidates(text)

    def keep(self, doc: RawDocument, cleaned_text: str) -> bool:
        if doc.source_type == "exercise":
            return bool(cleaned_text)
        return self._relevance_score(cleaned_text) >= 2

    def keep_chunk(self, doc: RawDocument, chunk_text: str) -> bool:
        if doc.source_type in {"exercise", "docx"}:
            return self._relevance_score(chunk_text) >= 1
        return self._relevance_score(chunk_text) >= 2

    def _relevance_score(self, text: str) -> int:
        lower_text = text.lower()
        matched_nodes = self.catalog.match_text(text)
        score = min(len(matched_nodes), 4) * 2
        if any(re.search(pattern, text, flags=re.IGNORECASE) for pattern in self.CODE_PATTERNS):
            score += 2
        if any(marker in lower_text for marker in self.LANGUAGE_MARKERS):
            score += 1
        if self._looks_like_background_page(lower_text) and score < 4:
            score -= 3
        return score

    @staticmethod
    def _looks_like_background_page(lower_text: str) -> bool:
        background = {
            "课程安排", "课程目标", "参考教材", "成绩评定", "发展历史", "作者", "版本历史", "参考文献",
        }
        return any(item in lower_text for item in background)


class ChunkSplitter:
    """
    知识片段切分器。

    为什么要切分？
    因为不能把整份 PPT 或整页 PDF 直接交给实体抽取。
    切成小片段后，每个片段更容易对应一个概念、一个语法规则或一个代码示例。
    """

    def split(self, doc: RawDocument, cleaned_text: str) -> list[str]:
        """把一段清洗后的文本切成若干个小 chunk。"""
        if doc.source_type == "exercise":
            # 习题天然就是一个相对完整的单元，所以不再拆得太碎。
            return [cleaned_text]

        blocks = self._split_by_blank_lines_or_bullets(cleaned_text)
        chunks: list[str] = []
        for block in blocks:
            if len(block) <= 420:
                chunks.append(block)
            else:
                # 如果某个段落太长，就继续按句子切分，避免后续 LLM 输入过长。
                chunks.extend(self._split_long_block(block))
        return [chunk.strip() for chunk in chunks if len(chunk.strip()) >= 8]

    @staticmethod
    def _split_by_blank_lines_or_bullets(text: str) -> list[str]:
        """优先按空行、项目符号、编号切分。PPT 内容通常很适合这样切。"""
        raw_blocks = re.split(r"\n\s*\n", text)
        blocks: list[str] = []
        for raw_block in raw_blocks:
            lines = [line.strip() for line in raw_block.splitlines() if line.strip()]
            current: list[str] = []
            for line in lines:
                is_bullet = bool(re.match(r"^([\-*•]|\d+[.)、]|[A-D][.)、])\s*", line))
                if is_bullet and current:
                    blocks.append(" ".join(current))
                    current = [line]
                else:
                    current.append(line)
            if current:
                blocks.append(" ".join(current))
        return blocks

    @staticmethod
    def _split_long_block(text: str) -> list[str]:
        """如果一个文本块太长，就按句号、问号、感叹号继续切分。"""
        sentences = re.split(r"(?<=[。！？.!?])\s+", text)
        chunks: list[str] = []
        current = ""
        for sentence in sentences:
            if not sentence:
                continue
            if len(current) + len(sentence) > 420 and current:
                chunks.append(current)
                current = sentence
            else:
                current = f"{current} {sentence}".strip()
        if current:
            chunks.append(current)
        return chunks


class ChunkLabeler:
    """
    轻量级规则标注器。

    它会给每个 chunk 补充一些标签，例如：
    - language：Java / Python / C++ / 通用OOP
    - material_role：概念解释 / 语法规则 / 代码示例 / 习题
    - chapter：大致属于哪个章节
    - normalized_content：归一化后的文本

    后续如果接入 LLM，可以优先替换这个类。
    只要输入输出结构不变，后面的流程不用大改。
    """

    def label(self, text: str, doc: RawDocument, keywords: list[str]) -> dict[str, Any]:
        """给一个 chunk 自动打标签。"""
        return {
            "language": self._detect_language(text, doc.source_file),
            "material_role": self._detect_role(text, doc.source_type),
            "chapter": self._detect_chapter(doc.title, text),
            "section": self._detect_section(doc.title, text, keywords),
            "normalized_content": self._normalize_terms(text),
            "confidence_hint": self._confidence_hint(text, keywords),
        }

    @staticmethod
    def _detect_language(text: str, source_file: str) -> str:
        """根据文件名和文本中的关键词，粗略判断编程语言。"""
        combined = f"{source_file}\n{text}".lower()
        if "java" in combined or "extends" in combined or "implements" in combined:
            return "Java"
        if "python" in combined or "__init__" in combined or "self." in combined:
            return "Python"
        if "c++" in combined or "virtual" in combined or "::" in combined:
            return "C++"
        return "通用编程"

    @staticmethod
    def _detect_role(text: str, source_type: str) -> str:
        """判断 chunk 是概念解释、语法规则、代码示例还是习题。"""
        if source_type == "exercise":
            return "exercise"
        lower_text = text.lower()
        scores: dict[str, int] = {}
        for role, role_keywords in ROLE_KEYWORDS.items():
            scores[role] = sum(1 for keyword in role_keywords if keyword.lower() in lower_text)
        best_role = max(scores, key=scores.get)
        if scores[best_role] == 0:
            return "concept_explanation"
        return best_role

    @staticmethod
    def _detect_chapter(title: str, text: str) -> str:
        """根据标题和正文内容，粗略归到某个面向对象主题下。"""
        combined = f"{title} {text}"
        if "继承" in combined or "extends" in combined:
            return "面向对象编程-继承"
        if "多态" in combined or "polymorphism" in combined.lower():
            return "面向对象编程-多态"
        if "接口" in combined or "interface" in combined.lower():
            return "面向对象编程-接口"
        if "封装" in combined or "private" in combined.lower() or "protected" in combined.lower():
            return "面向对象编程-封装"
        if "构造" in combined or "constructor" in combined.lower():
            return "面向对象编程-构造方法"
        if "类" in combined or "对象" in combined or "class" in combined.lower():
            return "面向对象编程-类与对象"
        return "面向对象编程"

    @staticmethod
    def _detect_section(title: str, text: str, keywords: list[str]) -> str:
        """小节名优先使用 PPT/PDF 的标题；没有标题时用关键词兜底。"""
        if title:
            return title
        return "、".join(keywords[:3]) if keywords else "未标注小节"

    @staticmethod
    def _normalize_terms(text: str) -> str:
        """把同义表达替换为统一术语，减少后续重复实体。"""
        normalized = text
        for old, new in NORMALIZATION_MAP.items():
            normalized = re.sub(re.escape(old), new, normalized, flags=re.IGNORECASE)
        return normalized

    @staticmethod
    def _confidence_hint(text: str, keywords: list[str]) -> float:
        """
        给一个简单的可信度提示分。

        这不是最终图谱质量评分，只是告诉后续流程：
        这个 chunk 看起来是否足够像 OOP 相关材料。
        """
        score = 0.55
        if len(keywords) >= 2:
            score += 0.15
        if any(mark in text for mark in ["class", "extends", "implements", "接口", "继承", "多态", "封装"]):
            score += 0.15
        if len(text) >= 25:
            score += 0.1
        return min(score, 0.95)


class ProgrammingChunkLabeler(ChunkLabeler):
    """编程领域标注器：章节归属来自标准目录，而不是固定 OOP 分支。"""

    def __init__(self, catalog: CurriculumCatalog) -> None:
        self.catalog = catalog

    def label(self, text: str, doc: RawDocument, keywords: list[str]) -> dict[str, Any]:
        domain_candidates = self.catalog.domain_candidates(f"{doc.title}\n{text}")
        labels = {
            "language": self._detect_language(text, doc.source_file),
            "material_role": self._detect_role(text, doc.source_type),
            "chapter": domain_candidates[0] if domain_candidates else "未归类编程内容",
            "section": self._detect_section(doc.title, text, keywords),
            "normalized_content": self._normalize_terms(text),
            "confidence_hint": min(0.95, 0.55 + 0.1 * min(len(domain_candidates), 3) + (0.1 if len(text) >= 25 else 0.0)),
            "domain_candidates": domain_candidates,
        }
        return labels


class MaterialPreprocessor:
    """
    总控类。

    它把前面的几个小模块串起来：

        读取文件 → 清洗文本 → 领域过滤 → 切分 chunk → 自动标注 → 输出 CleanChunk
    """

    def __init__(self, mode: str = "programming", catalog_path: str | Path = DEFAULT_CATALOG) -> None:
        # 每个成员变量都是流水线中的一个步骤。
        self.mode = mode
        self.loader = MaterialLoader()
        self.cleaner = TextCleaner()
        self.splitter = ChunkSplitter()
        if mode == "oop":
            self.filter = DomainFilter()
            self.labeler = ChunkLabeler()
        else:
            catalog = CurriculumCatalog.load(catalog_path)
            self.filter = ProgrammingDomainFilter(catalog)
            self.labeler = ProgrammingChunkLabeler(catalog)

    def run(self, input_dir: Path) -> list[CleanChunk]:
        """运行完整预处理流程，返回所有清洗后的知识片段。"""
        # 第一步：从输入目录读取所有支持的文件。
        raw_docs = self.loader.load_folder(input_dir)
        chunks: list[CleanChunk] = []

        for doc in raw_docs:
            # 第二步：清洗原始文本。
            cleaned_text = self.cleaner.clean(doc.raw_text)

            # 第三步：过滤掉与当前知识域无关的内容。
            if not cleaned_text or not self.filter.keep(doc, cleaned_text):
                continue

            # 第四步：把清洗后的文本切成更小的知识片段。
            for index, content in enumerate(self.splitter.split(doc, cleaned_text), start=1):
                keywords = self.filter.matched_keywords(content)
                if not keywords:
                    continue
                # 第四点五步：对单个 chunk 再做一次更严格的过滤。
                # 这样可以避免整页因为某个 OOP 词被保留后，页内的历史、版本、目录片段也混进来。
                if not self.filter.keep_chunk(doc, content):
                    continue

                # 第五步：给 chunk 自动补充语言、材料类型、章节等标签。
                labels = self.labeler.label(content, doc, keywords)
                chunk_id = self._make_chunk_id(doc, index, content)

                # 第六步：组装成统一输出结构。
                chunks.append(
                    CleanChunk(
                        chunk_id=chunk_id,
                        source_file=doc.source_file,
                        source_type=doc.source_type,
                        language=labels["language"],
                        page=doc.page,
                        chapter=labels["chapter"],
                        section=labels["section"],
                        content=content,
                        normalized_content=labels["normalized_content"],
                        keywords=keywords,
                        domain_candidates=labels.get("domain_candidates", []),
                        material_role=labels["material_role"],
                        evidence_location=self._evidence_location(doc),
                        confidence_hint=labels["confidence_hint"],
                        metadata=doc.metadata,
                    )
                )
        return chunks

    @staticmethod
    def _make_chunk_id(doc: RawDocument, index: int, content: str) -> str:
        """生成稳定的 chunk_id，方便后续追踪和去重。"""
        location = doc.question_id or doc.page or "doc"
        digest = hashlib.md5(f"{doc.source_file}-{location}-{index}-{content}".encode("utf-8")).hexdigest()[:8]
        source_stem = re.sub(r"[^a-zA-Z0-9]+", "_", Path(doc.source_file).stem).strip("_").lower()
        return f"{source_stem}_{location}_{index}_{digest}"

    @staticmethod
    def _evidence_location(doc: RawDocument) -> str:
        """生成证据位置，例如第 12 页、第 3 题。"""
        if doc.page is not None:
            return f"第{doc.page}页"
        if doc.question_id:
            return f"第{doc.question_id}题"
        return "全文"


def write_chunks(chunks: list[CleanChunk], output_path: Path) -> None:
    """把 CleanChunk 列表写入 clean_chunks.json。"""
    output_path.parent.mkdir(parents=True, exist_ok=True)
    payload = [chunk.__dict__ for chunk in chunks]
    output_path.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def parse_args(argv: list[str]) -> argparse.Namespace:
    """解析命令行参数。"""
    parser = argparse.ArgumentParser(description="Preprocess programming course materials into clean_chunks.json")
    parser.add_argument("--input", required=True, help="Folder containing PPT/PDF/DOCX/TXT/MD/JSON materials")
    parser.add_argument("--output", required=True, help="Output clean_chunks.json path")
    parser.add_argument("--mode", choices=["programming", "oop"], default="programming", help="programming 为正式模式；oop 仅用于复现旧版")
    parser.add_argument("--catalog", default=DEFAULT_CATALOG, help="编程领域标准知识目录 JSON 路径")
    return parser.parse_args(argv)


def main(argv: list[str]) -> int:
    """程序入口：读取参数，执行预处理，输出结果。"""
    args = parse_args(argv)
    input_dir = Path(args.input)
    output_path = Path(args.output)

    if not input_dir.exists():
        print(f"输入目录不存在：{input_dir}", file=sys.stderr)
        return 1

    preprocessor = MaterialPreprocessor(mode=args.mode, catalog_path=args.catalog)
    chunks = preprocessor.run(input_dir)
    write_chunks(chunks, output_path)
    print(f"已生成 {len(chunks)} 个知识片段：{output_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))
